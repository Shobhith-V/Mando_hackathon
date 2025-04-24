import fitz  # PyMuPDF
import io
from PIL import Image
from backend.ocr_engine import ocr_image_bytes  
from pptx import Presentation
from io import BytesIO
import os

# Parsing a PDF file by extracting text and running OCR on images
def parse_pdf(file):
    try:
        doc = fitz.open(stream=file.read(), filetype="pdf")
        full_text = []
        
        for page_num, page in enumerate(doc):
            # Extracting visible text from the current page
            page_text = page.get_text()
            if page_text:
                full_text.append(f"Page {page_num+1} Text:\n{page_text}")
            
            # Getting all images from the current page
            img_list = page.get_images(full=True)
            
            for img_index, img in enumerate(img_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                
                # Running OCR on each extracted image using a centralized function
                ocr_result = ocr_image_bytes(base_image["image"])
                
                full_text.append(
                    f"\nPage {page_num+1} Image {img_index+1} OCR:\n{ocr_result}"
                )
        
        return "\n".join(full_text)
    except Exception as e:
        return f"PDF Error: {str(e)}"
    

# Parsing a plain text file by decoding it to UTF-8
def parse_txt(file):
    try:
        return file.read().decode("utf-8")
    except Exception as e:
        return f"Error parsing TXT: {e}"

import json

# Parsing a JSON file and returning a dictionary or list
def parse_json(file):
    try:
        content = file.read().decode("utf-8")
        data = json.loads(content)
        return data
    except Exception as e:
        return f"JSON Error: {str(e)}"

import pandas as pd

# Reading a CSV into a DataFrame for structured analysis
def parse_csv(file):
    try:
        df = pd.read_csv(file)
        return df
    except Exception as e:
        return f"CSV Error: {str(e)}"

# Reading an Excel file into a DataFrame
def parse_xlsx(file):
    try:
        df = pd.read_excel(file)
        return df
    except Exception as e:
        return f"XLSX Error: {str(e)}"

# Parsing a PPTX file by extracting text and saving embedded images
def parse_pptx(file, image_output_dir="pptx_images"):
    try:
        prs = Presentation(file)
        text_runs = []
        images = []

        # Creating output directory if it does not exist
        if not os.path.exists(image_output_dir):
            os.makedirs(image_output_dir)

        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
                if shape.shape_type == 13:  # Checking if the shape is an image (PICTURE)
                    image = shape.image
                    image_bytes = image.blob
                    image_format = image.ext  # Getting image format (e.g., 'jpeg', 'png')
                    image_filename = f"slide{slide_idx+1}_img{shape_idx+1}.{image_format}"
                    image_path = os.path.join(image_output_dir, image_filename)

                    # Writing image bytes to a file
                    with open(image_path, "wb") as f:
                        f.write(image_bytes)
                    images.append(image_path)

        result = {
            "text": "\n".join(text_runs),
            "images": images  # Returning file paths of extracted images
        }
        return result
    except Exception as e:
        return f"PPTX Error: {str(e)}"

# Routing to the appropriate file parser based on file type
def parse_file(file, filetype):
    try:
        if filetype == "pdf":
            return parse_pdf(file)
        elif filetype == "txt":
            return parse_txt(file)
        elif filetype == "json":
            return parse_json(file)
        elif filetype == "csv":
            return parse_csv(file)
        elif filetype == "xlsx":
            return parse_xlsx(file)
        elif filetype == "pptx":
            return parse_pptx(file)
        else:
            return "Unsupported file type"
    except Exception as e:
        return f"Error parsing file: {e}"
